#!/usr/bin/env python
# -*- coding: utf-8 -*-

import os
import re
import pandas as pd
import numpy as np
import torch
import torch.nn as nn
import torch.nn.functional as F
from torch_geometric.nn import SAGEConv
from py2neo import Graph, Node, Relationship
import jieba
import jieba.analyse
from pdfminer.high_level import extract_text
from pptx import Presentation
import subprocess
import json
from sklearn.metrics import f1_score, accuracy_score

# 第一部分：数据采集与预处理
class DataCollector:
    def __init__(self, textbook_path, ppt_dir, syllabus_path):
        self.textbook_path = textbook_path
        self.ppt_dir = ppt_dir
        self.syllabus_path = syllabus_path
        
    def extract_from_textbook(self):
        """从教材PDF中提取文本"""
        print("从教材中提取知识点...")
        text = extract_text(self.textbook_path)
        # 提取知识点定义到案例/应用之间的内容
        knowledge_points = []
        # 使用正则表达式匹配知识点定义到案例之间的内容
        pattern = r'(.*?定义.*?)(?=案例|应用|例子|$)'
        matches = re.finditer(pattern, text, re.DOTALL)
        for match in matches:
            knowledge_points.append(match.group(1).strip())
        return knowledge_points
    
    def extract_from_ppt(self):
        """从PPT课件中提取知识点"""
        print("从PPT课件中提取知识点...")
        knowledge_points = []
        for ppt_file in os.listdir(self.ppt_dir):
            if ppt_file.endswith('.pptx'):
                ppt_path = os.path.join(self.ppt_dir, ppt_file)
                prs = Presentation(ppt_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text = shape.text
                            if text and ("定义" in text or "概念" in text):
                                knowledge_points.append(text.strip())
        return knowledge_points
    
    def extract_from_syllabus(self):
        """从教学大纲Excel表中提取知识点"""
        print("从教学大纲中提取知识点...")
        df = pd.read_excel(self.syllabus_path)
        # 假设教学大纲中有一列叫做"知识点"
        if "知识点" in df.columns:
            return df["知识点"].dropna().tolist()
        else:
            # 尝试从其他列中提取
            knowledge_points = []
            for col in df.columns:
                for cell in df[col]:
                    if isinstance(cell, str) and ("定义" in cell or "概念" in cell):
                        knowledge_points.append(cell.strip())
            return knowledge_points
    
    def extract_from_video(self, video_path, api_key):
        """从教学视频中提取音频并转换为文本"""
        print("从教学视频中提取音频并转换为文本...")
        # 使用FFmpeg提取音频
        audio_path = video_path.replace('.mp4', '.wav')
        subprocess.call(['ffmpeg', '-i', video_path, '-vn', '-acodec', 'pcm_s16le', '-ar', '16000', '-ac', '1', audio_path])
        
        # 使用科大讯飞API进行语音转文字
        # 这里只是示例，实际使用需要根据科大讯飞API的具体调用方式
        # 假设返回的结果包含时间戳和文本内容
        # result = xunfei_api_call(audio_path, api_key)
        
        # 模拟API返回结果
        result = [
            {"timestamp": "00:01:15", "text": "今天我们讲解人工智能的基本概念..."},
            {"timestamp": "00:05:30", "text": "神经网络是一种模拟人脑结构和功能的数学模型..."}
        ]
        
        return result
    
    def collect_all_data(self):
        """收集所有数据源的知识点"""
        textbook_data = self.extract_from_textbook()
        ppt_data = self.extract_from_ppt()
        syllabus_data = self.extract_from_syllabus()
        
        # 合并所有数据
        all_data = textbook_data + ppt_data + syllabus_data
        print(f"共收集到{len(all_data)}个原始知识点")
        return all_data

# 第二部分：数据清洗与规范化
class DataCleaner:
    def __init__(self):
        self.min_length = 100
        self.max_length = 1000
        
    def clean_text(self, text):
        """清洗文本，去除标点符号等"""
        # 去除多余的空白字符
        text = re.sub(r'\s+', ' ', text).strip()
        # 去除特殊字符和标点符号（保留中文标点）
        text = re.sub(r'[^\w\s\u4e00-\u9fff，。！？；：""''（）【】《》]', '', text)
        return text
    
    def filter_by_length(self, texts):
        """过滤过长或过短的文本"""
        return [text for text in texts if self.min_length <= len(text) <= self.max_length]
    
    def remove_duplicates(self, texts):
        """去除重复的文本"""
        return list(set(texts))
    
    def handle_missing_values(self, texts):
        """处理缺失值"""
        return [text for text in texts if text and not text.isspace()]
    
    def process(self, texts):
        """处理文本数据"""
        print("开始清洗数据...")
        # 清洗每个文本
        cleaned_texts = [self.clean_text(text) for text in texts]
        # 处理缺失值
        cleaned_texts = self.handle_missing_values(cleaned_texts)
        # 去除重复
        cleaned_texts = self.remove_duplicates(cleaned_texts)
        # 过滤长度
        cleaned_texts = self.filter_by_length(cleaned_texts)
        
        print(f"清洗后剩余{len(cleaned_texts)}个知识点")
        return cleaned_texts

# 第三部分：知识图谱构建与图神经网络模型
class KnowledgeGraphBuilder:
    def __init__(self, uri="bolt://localhost:7687", user="neo4j", password="password"):
        self.graph = Graph(uri, auth=(user, password))
        
    def create_knowledge_graph(self, knowledge_points):
        """创建知识图谱"""
        print("开始构建知识图谱...")
        # 清空数据库
        self.graph.run("MATCH (n) DETACH DELETE n")
        
        # 创建课程节点
        course_node = Node("Course", name="大数据技术")
        self.graph.create(course_node)
        
        # 创建知识点节点
        for i, kp in enumerate(knowledge_points):
            # 提取关键词作为知识点名称
            keywords = jieba.analyse.extract_tags(kp, topK=3)
            kp_name = "".join(keywords) if keywords else f"知识点{i}"
            
            # 创建知识点节点
            kp_node = Node("KnowledgeNode", name=kp_name, description=kp)
            self.graph.create(kp_node)
            
            # 创建与课程的关系
            rel = Relationship(course_node, "CONTAINS", kp_node)
            self.graph.create(rel)
        
        # 建立知识点之间的关系
        self._build_relationships()
        
        print("知识图谱构建完成")
        
    def _build_relationships(self):
        """建立知识点之间的关系"""
        # 获取所有知识点
        query = "MATCH (k:KnowledgeNode) RETURN k"
        knowledge_nodes = self.graph.run(query).data()
        
        # 建立关系
        for i, node1 in enumerate(knowledge_nodes):
            for j, node2 in enumerate(knowledge_nodes):
                if i != j:
                    # 计算两个知识点的相似度（简化示例）
                    similarity = self._calculate_similarity(
                        node1['k']['description'], 
                        node2['k']['description']
                    )
                    
                    # 如果相似度高，建立RELATED_TO关系
                    if similarity > 0.5:
                        query = """
                        MATCH (k1:KnowledgeNode), (k2:KnowledgeNode)
                        WHERE id(k1) = $id1 AND id(k2) = $id2
                        CREATE (k1)-[r:RELATED_TO {weight: $weight}]->(k2)
                        """
                        self.graph.run(query, id1=node1['k'].identity, 
                                      id2=node2['k'].identity, 
                                      weight=similarity)
    
    def _calculate_similarity(self, text1, text2):
        """计算两段文本的相似度"""
        # 这里使用简化的方法，实际应用中可以使用更复杂的算法
        words1 = set(jieba.cut(text1))
        words2 = set(jieba.cut(text2))
        
        # 计算Jaccard相似度
        intersection = len(words1.intersection(words2))
        union = len(words1.union(words2))
        
        return intersection / union if union > 0 else 0
    
    def export_graph_data(self):
        """导出图数据，用于图神经网络训练"""
        # 获取所有节点
        nodes_query = "MATCH (k:KnowledgeNode) RETURN id(k) AS id, k.description AS text"
        nodes = self.graph.run(nodes_query).data()
        
        # 获取所有边
        edges_query = """
        MATCH (k1:KnowledgeNode)-[r:RELATED_TO]->(k2:KnowledgeNode)
        RETURN id(k1) AS source, id(k2) AS target, r.weight AS weight
        """
        edges = self.graph.run(edges_query).data()
        
        # 转换为适合图神经网络的格式
        node_features = []
        for node in nodes:
            # 这里简化处理，实际应用中需要将文本转换为特征向量
            node_features.append({"id": node["id"], "text": node["text"]})
        
        edge_index = []
        edge_weight = []
        for edge in edges:
            edge_index.append([edge["source"], edge["target"]])
            edge_weight.append(edge["weight"])
        
        return {
            "nodes": node_features,
            "edge_index": edge_index,
            "edge_weight": edge_weight
        }

# GraphSAGE模型定义
class GraphSAGE(nn.Module):
    def __init__(self, in_channels, hidden_channels, out_channels):
        super(GraphSAGE, self).__init__()
        self.conv1 = SAGEConv(in_channels, hidden_channels)
        self.conv2 = SAGEConv(hidden_channels, out_channels)
        
    def forward(self, x, edge_index):
        x = self.conv1(x, edge_index)
        x = F.relu(x)
        x = F.dropout(x, p=0.5, training=self.training)
        x = self.conv2(x, edge_index)
        return x

# 第四部分：模型评估
class ModelEvaluator:
    def __init__(self):
        pass
    
    def evaluate_entity_extraction(self, true_entities, pred_entities):
        """评估实体识别性能"""
        # 计算准确率、召回率、F1值
        # 这里简化处理，实际应用需要更复杂的评估方法
        true_set = set(true_entities)
        pred_set = set(pred_entities)
        
        precision = len(true_set.intersection(pred_set)) / len(pred_set) if pred_set else 0
        recall = len(true_set.intersection(pred_set)) / len(true_set) if true_set else 0
        f1 = 2 * precision * recall / (precision + recall) if (precision + recall) > 0 else 0
        
        print(f"实体识别 F1值: {f1:.4f}")
        return f1
    
    def evaluate_relation_extraction(self, true_relations, pred_relations):
        """评估关系提取性能"""
        # 假设关系是(entity1, relation_type, entity2)的三元组
        true_set = set(true_relations)
        pred_set = set(pred_relations)
        
        accuracy = len(true_set.intersection(pred_set)) / len(pred_set) if pred_set else 0
        print(f"关系识别准确率: {accuracy:.4f}")
        return accuracy
    
    def evaluate_recommendation(self, expert_scores, student_scores):
        """评估知识推荐质量"""
        # 计算平均满意度
        all_scores = expert_scores + student_scores
        avg_satisfaction = sum(all_scores) / len(all_scores) if all_scores else 0
        print(f"推荐满意度: {avg_satisfaction:.4f}")
        return avg_satisfaction
    
    def ab_testing(self, traditional_scores, kg_based_scores):
        """A/B测试比较传统学习与基于知识图谱的学习效果"""
        # 计算平均分数
        avg_traditional = sum(traditional_scores) / len(traditional_scores)
        avg_kg_based = sum(kg_based_scores) / len(kg_based_scores)
        
        # 计算提升百分比
        improvement = (avg_kg_based - avg_traditional) / avg_traditional * 100
        print(f"学习效果提升: {improvement:.2f}%")
        return improvement

# 主函数
def main():
    # 设置路径
    textbook_path = "data/textbook.pdf"
    ppt_dir = "data/ppt"
    syllabus_path = "data/syllabus.xlsx"
    
    # 数据收集
    collector = DataCollector(textbook_path, ppt_dir, syllabus_path)
    all_data = collector.collect_all_data()
    
    # 数据清洗
    cleaner = DataCleaner()
    cleaned_data = cleaner.process(all_data)
    
    # 知识图谱构建
    kg_builder = KnowledgeGraphBuilder()
    kg_builder.create_knowledge_graph(cleaned_data)
    
    # 导出图数据
    graph_data = kg_builder.export_graph_data()
    
    # 构建和训练GraphSAGE模型（简化示例）
    # 实际应用中需要更完整的模型训练流程
    # 包括特征向量化、超参数调优等
    
    # 评估系统性能
    evaluator = ModelEvaluator()
    print("系统构建完成！大数据技术课程知识图谱已生成。")

if __name__ == "__main__":
    main() 